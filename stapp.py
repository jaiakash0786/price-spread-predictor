import streamlit as st
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from sklearn.ensemble import RandomForestRegressor
from sklearn.model_selection import train_test_split
from sklearn.metrics import mean_absolute_error, mean_squared_error, r2_score
import joblib
import os
from datetime import datetime
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches

st.set_page_config(page_title="Product Price Spread Predictor", layout="wide")

PRODUCT_LIST = [
    "Strawberries", "Romaine Lettuce", "Red Leaf Lettuce", "Potatoes", "Oranges",
    "Iceberg Lettuce", "Green Leaf Lettuce", "Celery", "Cauliflower", "Carrots",
    "Cantaloupe", "Broccoli Crowns", "Avocados"
]

# ------------------ LOAD DATA ------------------ #
@st.cache_data
def load_data():
    df = pd.read_csv("product_prices.csv")
    price_cols = ["farmprice", "atlantaretail", "chicagoretail", "losangelesretail", "newyorkretail"]
    for col in price_cols:
        df[col] = df[col].replace(r'[\$,]', '', regex=True).replace('', np.nan)
        df[col] = pd.to_numeric(df[col], errors='coerce')
    df["averagespread"] = df["averagespread"].str.replace('%', '', regex=False).replace('', np.nan)
    df["averagespread"] = pd.to_numeric(df["averagespread"], errors='coerce') / 100
    df.dropna(subset=price_cols + ["averagespread"], inplace=True)
    df["profit_atlanta"] = df["atlantaretail"] - df["farmprice"]
    df["profit_chicago"] = df["chicagoretail"] - df["farmprice"]
    df["profit_losangeles"] = df["losangelesretail"] - df["farmprice"]
    df["profit_newyork"] = df["newyorkretail"] - df["farmprice"]
    return df

# ------------------ TRAIN MODEL ------------------ #
@st.cache_resource
def train_model(df):
    X = df[["farmprice", "profit_atlanta", "profit_chicago", "profit_losangeles", "profit_newyork"]]
    y = df["averagespread"]
    X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.3, random_state=42)
    model = RandomForestRegressor(n_estimators=100, random_state=42)
    model.fit(X_train, y_train)
    joblib.dump(model, "rf_price_predictor.pkl")
    y_pred = model.predict(X_test)
    return model, y_test, y_pred, model.score(X_test, y_test)

# ------------------ PLOTTING FUNCTIONS ------------------ #
def plot_profit_vs_quantity(quantity, best_qty, predicted_spread, profit_unit, product, city):
    quantities = list(range(1, quantity + 1))
    profits = [q * profit_unit * predicted_spread for q in quantities]
    plt.figure(figsize=(8, 5))
    plt.plot(quantities, profits, color='green')
    plt.axvline(best_qty, color='red', linestyle='--', label=f'Best Quantity = {best_qty}')
    plt.title(f"Profit vs Quantity for {product} in {city}")
    plt.xlabel("Quantity")
    plt.ylabel("Expected Profit ($)")
    plt.legend()
    plt.grid(True)
    buf = BytesIO()
    plt.savefig(buf, format='png')
    plt.close()
    buf.seek(0)
    return buf

def plot_loss_vs_quantity(quantity, best_qty, predicted_spread, profit_unit, product, city):
    quantities = list(range(1, quantity + 1))
    losses_full = [(quantity - q) * profit_unit * (1 - predicted_spread) for q in quantities]
    remaining_stock = [quantity - q for q in quantities]
    plt.figure(figsize=(8, 5))
    plt.plot(quantities, losses_full, color='orange', label='Loss (Full Quantity)')
    plt.plot(quantities, remaining_stock, color='purple', label='Remaining Stock')
    plt.title(f"Loss & Remaining Stock vs Quantity for {product} in {city}")
    plt.xlabel("Quantity Purchased")
    plt.ylabel("Loss / Remaining Stock")
    plt.legend()
    plt.grid(True)
    buf = BytesIO()
    plt.savefig(buf, format='png')
    plt.close()
    buf.seek(0)
    return buf

def plot_export_city_profit(cities, profits):
    plt.figure(figsize=(8,5))
    plt.bar(cities, profits, color=['orange','blue','green','purple'])
    plt.title("Estimated Profit for Export by City")
    plt.xlabel("City")
    plt.ylabel("Profit ($)")
    plt.grid(True, axis='y')
    buf = BytesIO()
    plt.savefig(buf, format='png')
    plt.close()
    buf.seek(0)
    return buf

# ------------------ PREDICTION FUNCTIONS ------------------ #
def predict_shop_owner(product, quantity, city, df, model):
    city_map = {
        "Atlanta": "profit_atlanta",
        "Chicago": "profit_chicago",
        "Los Angeles": "profit_losangeles",
        "New York": "profit_newyork"
    }
    product_data = df[df["productname"].str.lower() == product.lower()]
    if product_data.empty:
        st.error("Product not found in dataset.")
        return None
    avg_farm = product_data["farmprice"].mean()
    avg_profits = [product_data[c].mean() for c in city_map.values()]
    features = ["farmprice", "profit_atlanta", "profit_chicago", "profit_losangeles", "profit_newyork"]
    input_df = pd.DataFrame([[avg_farm] + avg_profits], columns=features)
    predicted_spread = model.predict(input_df)[0]
    profit_unit = product_data[city_map[city]].mean()
    best_qty = int((quantity * predicted_spread) / (predicted_spread + 1))
    max_profit = best_qty * profit_unit * predicted_spread
    loss = (quantity - best_qty) * profit_unit * (1 - predicted_spread)
    return {
        "product": product, "city": city, "quantity": quantity,
        "predicted_spread": predicted_spread, "best_qty": best_qty,
        "max_profit": max_profit, "loss": loss, "profit_unit": profit_unit
    }

def predict_farmer_export(product, quantity, df, model):
    product_data = df[df["productname"].str.lower() == product.lower()]
    if product_data.empty:
        st.error("Product not found in dataset.")
        return None
    avg_farm = product_data["farmprice"].mean()
    avg_profits = [product_data[c].mean() for c in ["profit_atlanta", "profit_chicago", "profit_losangeles", "profit_newyork"]]
    features = ["farmprice", "profit_atlanta", "profit_chicago", "profit_losangeles", "profit_newyork"]
    input_df = pd.DataFrame([[avg_farm] + avg_profits], columns=features)
    predicted_spread = model.predict(input_df)[0]
    profits = [quantity * predicted_spread * avg for avg in avg_profits]
    cities = ["Atlanta", "Chicago", "Los Angeles", "New York"]
    best_city = cities[np.argmax(profits)]
    max_profit = max(profits)
    return {
        "product": product, "predicted_spread": predicted_spread,
        "cities": cities, "profits": profits, "best_city": best_city,
        "max_profit": max_profit
    }

# ------------------ PPTX FUNCTION ------------------ #
def create_pptx_shop(shop_result, shop_graphs):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = f"Shop Owner Prediction - {shop_result['product']} in {shop_result['city']}"
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(2))
    tf = textbox.text_frame
    tf.text = f"Predicted Spread: {shop_result['predicted_spread']:.2%}"
    tf.add_paragraph().text = f"Suggested Quantity: {shop_result['best_qty']} out of {shop_result['quantity']}"
    tf.add_paragraph().text = f"Expected Maximum Profit: ${shop_result['max_profit']:.2f}"
    tf.add_paragraph().text = f"Loss if Buying Full Quantity: ${shop_result['loss']:.2f}"
    slide.shapes.add_picture(shop_graphs[0], Inches(0.5), Inches(3.5), Inches(4.5), Inches(3))
    slide.shapes.add_picture(shop_graphs[1], Inches(5), Inches(3.5), Inches(4.5), Inches(3))
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf

def create_pptx_export(export_result, export_graph):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = f"Farmer Export Prediction - {export_result['product']}"
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(2))
    tf = textbox.text_frame
    tf.text = f"Predicted Spread: {export_result['predicted_spread']:.2%}"
    tf.add_paragraph().text = f"Best City to Export: {export_result['best_city']}"
    tf.add_paragraph().text = f"Maximum Expected Profit: ${export_result['max_profit']:.2f}"
    slide.shapes.add_picture(export_graph, Inches(1), Inches(3.5), Inches(8), Inches(4))
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf

# ------------------ MAIN APP ------------------ #
st.title("Product Price Spread Predictor")

df = load_data()
model, y_test, y_pred, accuracy = train_model(df)

st.subheader("Model Evaluation Metrics")
st.write(f"Accuracy: {accuracy:.2%}")
st.write(f"MAE: {mean_absolute_error(y_test, y_pred):.4f}")
st.write(f"MSE: {mean_squared_error(y_test, y_pred):.4f}")
st.write(f"RMSE: {np.sqrt(mean_squared_error(y_test, y_pred)):.4f}")
st.write(f"R² Score: {r2_score(y_test, y_pred):.2%}")
st.markdown("---")

# --- Shop Owner Prediction ---
st.subheader("Shop Owner Prediction")
product_shop = st.selectbox("Select Product", PRODUCT_LIST, key="shop")
city_shop = st.selectbox("Select City", ["Atlanta", "Chicago", "Los Angeles", "New York"], key="shop_city")
quantity_shop = st.number_input("Enter Quantity", min_value=1, step=1, key="shop_qty")

if st.button("Predict Shop Owner"):
    shop_result = predict_shop_owner(product_shop, quantity_shop, city_shop, df, model)
    if shop_result:
        st.write("### Inputs")
        st.write(f"Product: {shop_result['product']}")
        st.write(f"City: {shop_result['city']}")
        st.write(f"Quantity: {shop_result['quantity']}")
        st.write("### Prediction Results")
        st.write(f"Predicted Spread: {shop_result['predicted_spread']:.2%}")
        st.write(f"Suggested Quantity: {shop_result['best_qty']}")
        st.write(f"Expected Max Profit: ${shop_result['max_profit']:.2f}")
        st.write(f"Loss if Buying Full Quantity: ${shop_result['loss']:.2f}")
        
        graph1 = plot_profit_vs_quantity(quantity_shop, shop_result['best_qty'], shop_result['predicted_spread'], shop_result['profit_unit'], product_shop, city_shop)
        graph2 = plot_loss_vs_quantity(quantity_shop, shop_result['best_qty'], shop_result['predicted_spread'], shop_result['profit_unit'], product_shop, city_shop)
        st.image([graph1, graph2], width=400)
        
        pptx_file = create_pptx_shop(shop_result, [graph1, graph2])
        st.download_button("Download Shop Owner Report PPTX", pptx_file, file_name="Shop_Owner_Report.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

# --- Farmer Export Prediction ---
st.subheader("Farmer Export Prediction")
product_farmer = st.selectbox("Select Product", PRODUCT_LIST, key="farmer")
quantity_farmer = st.number_input("Enter Quantity", min_value=1, step=1, key="farmer_qty")

if st.button("Predict Farmer Export"):
    export_result = predict_farmer_export(product_farmer, quantity_farmer, df, model)
    if export_result:
        st.write("### Inputs")
        st.write(f"Product: {export_result['product']}")
        st.write(f"Quantity: {quantity_farmer}")
        st.write("### Prediction Results")
        st.write(f"Predicted Spread: {export_result['predicted_spread']:.2%}")
        st.write(f"Best City to Export: {export_result['best_city']}")
        st.write(f"Maximum Expected Profit: ${export_result['max_profit']:.2f}")
        
        graph_export = plot_export_city_profit(export_result['cities'], export_result['profits'])
        st.image(graph_export, width=600)
        
        pptx_file_export = create_pptx_export(export_result, graph_export)
        st.download_button("Download Farmer Export Report PPTX", pptx_file_export, file_name="Farmer_Export_Report.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
